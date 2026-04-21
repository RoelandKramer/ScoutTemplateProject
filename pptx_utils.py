"""Utility functions for filling FC Den Bosch scouting PowerPoint templates."""

import copy
import io
import os
import re
from lxml import etree
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.oxml.ns import qn
from pptx.shapes.picture import Movie
from pptx.util import Pt

YELLOW = RGBColor(0xFF, 0xD9, 0x32)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)

# Maximum dimension (EMU) a shape can have to be considered a star.
# Real stars are ~3–5 mm; decorative bars/ovals are much larger.
_MAX_STAR_DIM = 600_000   # ≈ 6.5 mm
_ROW_TOLERANCE = 200_000  # ≈ 2.2 mm  — vertical grouping tolerance


# ─── Template definitions ────────────────────────────────────────────────────
# Each position has NL + ENG variable names, weights, and per-club/language
# file variants with the correct slide indices.

_BASE = "Template2.0/Spelersprofielen Keynote - Powerpoint"

CLUBS = ["FC Den Bosch", "Pro Vercelli"]
CLUB_LANGUAGES = {
    "FC Den Bosch": ["NL", "ENG"],
    "Pro Vercelli": ["ENG"],
}

TEMPLATES = {
    "Goalkeeper": {
        "variables_nl":  ["Balbehandeling", "Voorkomen vs verdedigen", "Bal verwerking", "Wendbaarheid", "Moed", "Onverstoorbaarheid"],
        "variables_eng": ["Handling", "Prevention vs defending", "Distribution", "Agility", "Bravery", "Consistency"],
        "weights": [1.1, 1.1, 1.0, 1.1, 1.1, 1.0],
        "variants": {
            ("FC Den Bosch", "NL"):  {"file": f"{_BASE}/#01 Keeper/#1 Goalkeeper Scouting FC Den Bosch (NL).pptx",  "rating_slide_idx": 3, "detail_slides": list(range(7, 13))},
            ("FC Den Bosch", "ENG"): {"file": f"{_BASE}/#01 Keeper/#1 Goalkeeper Scouting FC Den Bosch (ENG).pptx", "rating_slide_idx": 3, "detail_slides": list(range(7, 13))},
            ("Pro Vercelli", "ENG"): {"file": f"{_BASE}/#01 Keeper/#1 Goalkeeper Scouting Pro Vercelli (ENG).pptx", "rating_slide_idx": 1, "detail_slides": list(range(5, 11))},
        },
    },
    "Wingback": {
        "variables_nl":  ["Overlap/underlap", "Voorzetten", "Dribbelen met bal", "Snelheid", "Uithoudingsvermogen", "Wendbaarheid", "Doorzettingsvermogen"],
        "variables_eng": ["Overlap/underlap", "Crosses", "Carries", "Speed", "Stamina", "Agility", "Determination"],
        "weights": [1.1, 1.0, 1.0, 1.1, 1.0, 1.0, 1.0],
        "variants": {
            ("FC Den Bosch", "NL"):  {"file": f"{_BASE}/#02 #05 Wingback/#2 #5 Wingbacks Scouting Den Bosch (NL).pptx",  "rating_slide_idx": 3, "detail_slides": list(range(7, 14))},
            ("FC Den Bosch", "ENG"): {"file": f"{_BASE}/#02 #05 Wingback/#2 #5 Wingbacks Scouting Den Bosch (ENG).pptx", "rating_slide_idx": 3, "detail_slides": list(range(7, 14))},
            ("Pro Vercelli", "ENG"): {"file": f"{_BASE}/#02 #05 Wingback/#2 #5 Wingbacks Scouting Pro Vercelli (ENG).pptx", "rating_slide_idx": 1, "detail_slides": list(range(5, 12))},
        },
    },
    "Centerback": {
        "variables_nl":  ["Positie kiezen (v)", "Defensieve Koppen", "Offensieve Passing", "Dribbelen met bal", "Duelkracht", "Snelheid", "Doorzettingsvermogen"],
        "variables_eng": ["Defensive positioning", "Defensive headers", "Offensive Passing", "Carries", "Powerplay", "Speed", "Determination"],
        "weights": [1.1, 1.1, 1.0, 1.0, 1.1, 1.1, 1.0],
        "variants": {
            ("FC Den Bosch", "NL"):  {"file": f"{_BASE}/#03 #04 Centerbacks/#3 #4 Centerbacks Scouting FC Den Bosch (NL).pptx",  "rating_slide_idx": 3, "detail_slides": list(range(7, 14))},
            ("FC Den Bosch", "ENG"): {"file": f"{_BASE}/#03 #04 Centerbacks/#3 #4 Centerbacks Scouting FC Den Bosch (ENG).pptx", "rating_slide_idx": 3, "detail_slides": list(range(7, 14))},
            ("Pro Vercelli", "ENG"): {"file": f"{_BASE}/#03 #04 Centerbacks/#3 #4 Centerbacks Scouting Pro Vercelli (ENG).pptx", "rating_slide_idx": 1, "detail_slides": list(range(5, 12))},
        },
    },
    "Deep Lying Playmaker": {
        "variables_nl":  ["Offensieve Passing", "Positie kiezen (v)", "Dribbelen met bal", "Positie kiezen (A)", "Duelkracht", "Uithoudingsvermogen", "Snelheid", "Doorzettingsvermogen", "Onverstoorbaarheid"],
        "variables_eng": ["Offensive Passing", "Defensive positioning", "Carries", "Offensive positioning", "Powerplay", "Stamina", "Speed", "Determination", "Consistency"],
        "weights": [1.1, 1.1, 1.0, 1.0, 1.1, 1.1, 1.0, 1.1, 1.0],
        "variants": {
            ("FC Den Bosch", "NL"):  {"file": f"{_BASE}/#06 Deep lying playmaker/#6 Deep lying playmaker Scouting FC Den Bosch (NL).pptx",  "rating_slide_idx": 3, "detail_slides": list(range(8, 17))},
            ("FC Den Bosch", "ENG"): {"file": f"{_BASE}/#06 Deep lying playmaker/#6 Deep lying playmaker Scouting FC Den Bosch (ENG).pptx", "rating_slide_idx": 3, "detail_slides": list(range(8, 17))},
            ("Pro Vercelli", "ENG"): {"file": f"{_BASE}/#06 Deep lying playmaker/#6 Deep lying playmaker Scouting Pro Vercelli (ENG).pptx", "rating_slide_idx": 1, "detail_slides": list(range(6, 15))},
        },
    },
    "Box-to-Box Midfielder": {
        "variables_nl":  ["Offensieve Passing", "Positie kiezen (v)", "Dribbelen met bal", "Voorbij spelen tegenstander", "Positie kiezen (A)", "Duelkracht", "Uithoudingsvermogen", "Snelheid", "Doorzettingsvermogen"],
        "variables_eng": ["Offensive Passing", "Defensive positioning", "Carries", "Dribbling past opponent", "Offensive positioning", "Powerplay", "Stamina", "Speed", "Determination"],
        "weights": [1.1, 1.1, 1.0, 1.0, 1.0, 1.1, 1.1, 1.0, 1.1],
        "variants": {
            ("FC Den Bosch", "NL"):  {"file": f"{_BASE}/#08 Box-to-Box midfielder/#8 Box-to-Box midfielder Scouting FC Den Bosch (NL).pptx",  "rating_slide_idx": 3, "detail_slides": list(range(8, 17))},
            ("FC Den Bosch", "ENG"): {"file": f"{_BASE}/#08 Box-to-Box midfielder/#8 Box-to-Box midfielder Scouting FC Den Bosch (ENG).pptx", "rating_slide_idx": 3, "detail_slides": list(range(8, 17))},
            ("Pro Vercelli", "ENG"): {"file": f"{_BASE}/#08 Box-to-Box midfielder/#8 Box-to-Box midfielder Scouting Pro Vercelli (ENG).pptx", "rating_slide_idx": 1, "detail_slides": list(range(6, 15))},
        },
    },
    "Scoring 10": {
        "variables_nl":  ["Doelgerichtheid", "Diepte loopacties", "Positie kiezen (A)", "Duelkracht", "Uithoudingsvermogen", "Doorzettingsvermogen"],
        "variables_eng": ["Composure", "Deep runs", "Offensive positioning", "Powerplay", "Stamina", "Determination"],
        "weights": [1.1, 1.0, 1.0, 1.1, 1.0, 1.1],
        "variants": {
            ("FC Den Bosch", "NL"):  {"file": f"{_BASE}/#10 Scoring/#10 Scoring 10 Scouting FC Den Bosch (NL).pptx",  "rating_slide_idx": 3, "detail_slides": list(range(7, 13))},
            ("FC Den Bosch", "ENG"): {"file": f"{_BASE}/#10 Scoring/#10 Scoring 10 Scouting FC Den Bosch (ENG).pptx", "rating_slide_idx": 3, "detail_slides": list(range(7, 13))},
            ("Pro Vercelli", "ENG"): {"file": f"{_BASE}/#10 Scoring/#10 Scoring 10 Scouting Pro Vercelli (ENG).pptx", "rating_slide_idx": 1, "detail_slides": list(range(5, 11))},
        },
    },
    "Dribbling Winger": {
        "variables_nl":  ["Voorbij spelen tegenstander", "Doelgerichtheid", "Voorzet", "Afstandsschot", "Wendbaarheid", "Snelheid", "Flair", "Moed"],
        "variables_eng": ["Dribbling past opponent", "Composure", "Passing", "Shooting from distance", "Agility", "Speed", "Flair", "Bravery"],
        "weights": [1.1, 1.1, 1.0, 1.0, 1.1, 1.0, 1.1, 1.0],
        "variants": {
            ("FC Den Bosch", "NL"):  {"file": f"{_BASE}/#07 #11 Dribbling winger/#07 #11 Dribbling winger Scouting FC Den Bosch (NL).pptx",  "rating_slide_idx": 3, "detail_slides": list(range(7, 15))},
            ("FC Den Bosch", "ENG"): {"file": f"{_BASE}/#07 #11 Dribbling winger/#07 #11 Dribbling winger Scouting FC Den Bosch (ENG).pptx", "rating_slide_idx": 3, "detail_slides": list(range(7, 15))},
            ("Pro Vercelli", "ENG"): {"file": f"{_BASE}/#07 #11 Dribbling winger/#07 #11 Dribbling winger Scouting Pro Vercelli (ENG).pptx", "rating_slide_idx": 1, "detail_slides": list(range(5, 13))},
        },
    },
    "Fast Winger": {
        "variables_nl":  ["Voorbij tegenstander", "Diepte loopacties", "Doelgericht", "Snelheid", "Wendbaarheid", "Moed", "Flair"],
        "variables_eng": ["Dribbling past opponent", "Deep runs", "Composure", "Speed", "Agility", "Bravery", "Flair"],
        "weights": [1.1, 1.1, 1.0, 1.1, 1.1, 1.1, 1.0],
        "variants": {
            ("FC Den Bosch", "NL"):  {"file": f"{_BASE}/#07 #11 Fast winger/#07 #11 Fast winger Scouting FC Den Bosch (NL).pptx",  "rating_slide_idx": 3, "detail_slides": list(range(7, 14))},
            ("FC Den Bosch", "ENG"): {"file": f"{_BASE}/#07 #11 Fast winger/#07 #11 Fast winger Scouting FC Den Bosch (ENG).pptx", "rating_slide_idx": 3, "detail_slides": list(range(7, 14))},
            ("Pro Vercelli", "ENG"): {"file": f"{_BASE}/#07 #11 Fast winger/#07 #11 Fast winger Scouting Pro Vercelli (ENG).pptx", "rating_slide_idx": 1, "detail_slides": list(range(5, 12))},
        },
    },
    "Finisher": {
        "variables_nl":  ["Doelgerichtheid", "Positie kiezen (A)", "Combinatie spel", "Offensief koppen", "Duelkracht", "Flair", "Doorzettingsvermogen"],
        "variables_eng": ["Composure", "Offensive positioning", "Link up play", "Offensive headers", "Powerplay", "Flair", "Determination"],
        "weights": [1.1, 1.1, 1.0, 1.0, 1.1, 1.1, 1.0],
        "variants": {
            ("FC Den Bosch", "NL"):  {"file": f"{_BASE}/#09 Finisher/#9 Finisher Scouting FC Den Bosch (NL).pptx",  "rating_slide_idx": 3, "detail_slides": list(range(7, 14))},
            ("FC Den Bosch", "ENG"): {"file": f"{_BASE}/#09 Finisher/#9 Finisher Scouting FC Den Bosch (ENG).pptx", "rating_slide_idx": 3, "detail_slides": list(range(7, 14))},
            ("Pro Vercelli", "ENG"): {"file": f"{_BASE}/#09 Finisher/#9 Finisher Scouting Pro Vercelli (ENG).pptx", "rating_slide_idx": 1, "detail_slides": list(range(5, 12))},
        },
    },
}


def get_template_config(position: str, club: str, language: str) -> dict:
    """Build a flat config dict for the given position + club + language combo."""
    tmpl = TEMPLATES[position]
    variant = tmpl["variants"][(club, language)]
    lang_key = "variables_nl" if language == "NL" else "variables_eng"
    return {
        "file": variant["file"],
        "rating_slide_idx": variant["rating_slide_idx"],
        "detail_slides": variant["detail_slides"],
        "variables": tmpl[lang_key],
        "weights": tmpl["weights"],
    }


# ─── Competency description extraction ──────────────────────────────────────

_desc_cache: dict[str, list[dict]] = {}   # keyed by template file path


def extract_competency_descriptions(template_cfg: dict) -> list[dict]:
    """Extract competency name, description and assessment criteria from the
    explanation slide(s) of a template.

    Templates with many variables (e.g. 9) may spread their competency
    explanations across multiple slides immediately before the first detail
    slide.  We scan every slide between the rating slide and the first detail
    slide, skipping those that lack description text boxes (like the field
    diagram and the calculation table).

    Each explanation slide uses three kinds of text boxes:
      - 'TextBox 17': competency **name**
      - 'TextBox 30': short **description**
      - 'TextBox 31': newline-separated **assessment criteria**

    Returns one dict per variable, same order as the template's variable list:
        [{"name": "Handling",
          "description": "Defending a lot of balls",
          "criteria": ["Inside 16m", "Outside 16m", ...]}, ...]
    """
    file_path = template_cfg["file"]
    if file_path in _desc_cache:
        return _desc_cache[file_path]

    detail_slides = template_cfg.get("detail_slides", [])
    if not detail_slides:
        return []

    try:
        from pptx import Presentation as _Prs
        prs = _Prs(file_path)
    except Exception:
        return []

    # Scan all slides between rating_slide and first detail_slide that
    # contain description text boxes (TextBox 30).
    first_detail = detail_slides[0]
    rating_idx = template_cfg.get("rating_slide_idx", 0)

    names_boxes = []   # (left, top, text)  — global across explanation slides
    desc_boxes = []
    crit_boxes = []

    for si in range(rating_idx + 1, first_detail):
        if si < 0 or si >= len(prs.slides):
            continue
        slide = prs.slides[si]
        has_descs = False
        for sh in slide.shapes:
            nm = sh.name or ""
            if nm.startswith("TextBox 30") or nm == "TextBox 30":
                has_descs = True
                break
        if not has_descs:
            continue  # skip field diagram / table slides

        # Use a large offset on `top` per slide so shapes from later slides
        # sort after those on earlier slides.
        slide_offset = si * 100 * 914400

        for sh in slide.shapes:
            if not sh.has_text_frame:
                continue
            txt = sh.text_frame.text.strip().strip("\v").strip("\r")
            if not txt:
                continue
            nm = sh.name or ""
            adj_top = sh.top + slide_offset
            if nm.startswith("TextBox 17") or nm == "TextBox 17":
                names_boxes.append((sh.left, adj_top, txt))
            elif nm.startswith("TextBox 30") or nm == "TextBox 30":
                desc_boxes.append((
                    sh.left, adj_top,
                    txt.strip("\u200b").strip('"').strip("\u201c\u201d").strip(),
                ))
            elif nm.startswith("TextBox 31") or nm == "TextBox 31":
                items = [
                    line.strip().strip("\u2022").strip()
                    for line in txt.split("\n") if line.strip()
                ]
                crit_boxes.append((sh.left, adj_top, items))

    variables = template_cfg["variables"]

    # Sort name boxes in reading order (top then left)
    names_boxes.sort(key=lambda b: (b[1], b[0]))

    def _column_match(box_left, ref_left, tolerance=2 * 914400):
        return abs(box_left - ref_left) < tolerance

    results = []
    matched_names: set[int] = set()

    for var_name in variables:
        var_lower = var_name.lower().strip()
        best = None
        # Exact match first
        for i, (left, top, txt) in enumerate(names_boxes):
            if i in matched_names:
                continue
            if txt.lower().strip() == var_lower:
                best = i
                break
        # Substring match
        if best is None:
            for i, (left, top, txt) in enumerate(names_boxes):
                if i in matched_names:
                    continue
                if var_lower in txt.lower() or txt.lower().strip() in var_lower:
                    best = i
                    break
        # Take next unmatched in order
        if best is None:
            for i in range(len(names_boxes)):
                if i not in matched_names:
                    best = i
                    break

        if best is None:
            results.append({"name": var_name, "description": "", "criteria": []})
            continue

        matched_names.add(best)
        n_left, n_top, n_text = names_boxes[best]

        # Find matching description (same column, closest below the name)
        desc_text = ""
        best_dist = float("inf")
        for d_left, d_top, d_txt in desc_boxes:
            if _column_match(d_left, n_left) and d_top > n_top:
                dist = d_top - n_top
                if dist < best_dist:
                    best_dist = dist
                    desc_text = d_txt

        # Find matching criteria (same column, closest below the name)
        crit_items: list[str] = []
        best_dist = float("inf")
        for c_left, c_top, c_items in crit_boxes:
            if _column_match(c_left, n_left) and c_top > n_top:
                dist = c_top - n_top
                if dist < best_dist:
                    best_dist = dist
                    crit_items = c_items

        results.append({
            "name": n_text,
            "description": desc_text,
            "criteria": crit_items,
        })

    _desc_cache[file_path] = results
    return results


# ─── Star detection ──────────────────────────────────────────────────────────

def _is_star_shape(shape) -> bool:
    """Return True if the shape is a star (auto_shape type 92 OR small FREEFORM)."""
    try:
        if shape.auto_shape_type == 92:
            return True
    except (ValueError, AttributeError):
        pass
    if shape.shape_type == 5:  # FREEFORM
        return shape.width <= _MAX_STAR_DIM and shape.height <= _MAX_STAR_DIM
    return False


def get_star_rows(slide) -> list[list]:
    """Return star rows sorted top-to-bottom, each row sorted left-to-right.

    Handles both auto_shape type-92 stars (old templates) and small FREEFORM
    stars (newer templates).  Rows are identified by grouping shapes whose
    vertical positions are within _ROW_TOLERANCE of each other.
    """
    stars = [s for s in slide.shapes if _is_star_shape(s)]
    if not stars:
        return []

    row_map: dict = {}
    for star in stars:
        matched = False
        for rep_top in list(row_map.keys()):
            if abs(star.top - rep_top) < _ROW_TOLERANCE:
                row_map[rep_top].append(star)
                matched = True
                break
        if not matched:
            row_map[star.top] = [star]

    # Only keep rows that have at least 5 stars (filter out decorative singletons)
    rows = {top: shapes for top, shapes in row_map.items() if len(shapes) >= 5}

    return [
        sorted(rows[top], key=lambda s: s.left)
        for top in sorted(rows.keys())
    ]


# ─── Core operations ─────────────────────────────────────────────────────────

# OOXML gradient XML for a sharp left-half-yellow / right-half-white split.
# ang="0" in OOXML means the gradient flows left-to-right (constant-colour
# lines are vertical), so pos=0 is the left edge and pos=100000 the right.
_HALF_STAR_GRAD_XML = (
    '<a:gradFill xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"'
    ' rotWithShape="1">'
    '<a:gsLst>'
    '<a:gs pos="0"><a:srgbClr val="FFD932"/></a:gs>'
    '<a:gs pos="49999"><a:srgbClr val="FFD932"/></a:gs>'
    '<a:gs pos="50000"><a:srgbClr val="FFFFFF"/></a:gs>'
    '<a:gs pos="100000"><a:srgbClr val="FFFFFF"/></a:gs>'
    '</a:gsLst>'
    '<a:lin ang="0" scaled="0"/>'
    '</a:gradFill>'
)


def _apply_half_star_fill(shape) -> None:
    """Replace the shape's fill with a sharp left-yellow / right-white gradient."""
    # Let python-pptx create the solidFill element in the correct position
    shape.fill.solid()
    shape.fill.fore_color.rgb = YELLOW  # placeholder; will be replaced below

    spPr = shape._element.spPr
    solid = spPr.find(qn('a:solidFill'))
    if solid is None:
        return
    idx = list(spPr).index(solid)
    spPr.remove(solid)
    spPr.insert(idx, etree.fromstring(_HALF_STAR_GRAD_XML))


def _stop_srgb(gs_element) -> str | None:
    """Return the uppercase srgbClr val from an <a:gs> stop, or None."""
    clr = gs_element.find(qn('a:srgbClr'))
    return clr.get('val', '').upper() if clr is not None else None


def _stop_is_white(gs_element) -> bool:
    """True if the gradient stop is explicitly srgbClr FFFFFF."""
    srgb = gs_element.find(qn('a:srgbClr'))
    return srgb is not None and srgb.get('val', '').upper() == 'FFFFFF'


def _stop_is_non_white(gs_element) -> bool:
    """True if the gradient stop carries a color that is definitely not white.

    Handles both srgbClr (must be non-FFFFFF) and schemeClr/sysClr (always
    treated as non-white; Keynote maps FFD932 yellow → schemeClr accent4).
    """
    srgb = gs_element.find(qn('a:srgbClr'))
    if srgb is not None:
        return srgb.get('val', '').upper() not in ('FFFFFF', '')
    # schemeClr, sysClr, etc. — if ANY color child is present, assume non-white
    return len(gs_element) > 0


def _get_star_fill_value(shape) -> float:
    """Return 1.0 (full yellow), 0.5 (half-yellow gradient), or 0.0 (empty).

    Handles all known fill representations including Keynote round-trips where
    srgbClr FFD932 is remapped to schemeClr accent4:

      • solidFill srgbClr FFD932            → 1.0
      • solidFill srgbClr FFFFFF            → 0.0
      • solidFill schemeClr/sysClr (any)    → 1.0  (non-white scheme color)
      • our half-star gradFill (4 stops,    → 0.5
          FFD932+FFD932 / FFFFFF+FFFFFF)
      • Keynote half-star gradFill (same    → 0.5
          positions, accent4 instead of FFD932)
      • gradFill all FFFFFF stops           → 0.0
      • any other gradFill                  → 1.0  (assume filled)
    """
    try:
        spPr = shape._element.spPr

        grad = spPr.find(qn('a:gradFill'))
        if grad is not None:
            gsLst = grad.find(qn('a:gsLst'))
            stops = gsLst.findall(qn('a:gs')) if gsLst is not None else []

            if len(stops) == 4:
                positions = [int(s.get('pos', '0')) for s in stops]
                # Both our original and Keynote-modified half-star use these
                # exact positions — only the color type changes (srgb vs scheme)
                if positions == [0, 49999, 50000, 100000]:
                    low_non_white  = all(_stop_is_non_white(stops[j]) for j in (0, 1))
                    high_white     = all(_stop_is_white(stops[j])     for j in (2, 3))
                    if low_non_white and high_white:
                        return 0.5

            # All stops explicitly white → empty star
            if stops and all(_stop_is_white(s) for s in stops):
                return 0.0
            # Any yellow srgb stop → Keynote solid-yellow stored as gradient
            if any(_stop_srgb(s) == 'FFD932' for s in stops):
                return 1.0
            # Other gradient (theme color, etc.) — assume filled
            if stops:
                return 1.0

        solid = spPr.find(qn('a:solidFill'))
        if solid is not None:
            clr = solid.find(qn('a:srgbClr'))
            if clr is not None:
                val = clr.get('val', '').upper()
                if val == 'FFD932':
                    return 1.0
                if val == 'FFFFFF':
                    return 0.0
                return 1.0  # any other explicit rgb → treat as filled
            # schemeClr / sysClr / etc. — empty stars are always explicit white
            if len(solid) > 0:
                return 1.0

    except Exception:
        pass
    return 0.0


def color_stars(slide, star_values: list) -> None:
    """Colour stars from each row according to values (supports 0.5 increments).

    For a value of 7.5: stars 0-6 are full yellow, star 7 is half yellow,
    stars 8-9 are white.
    """
    rows = get_star_rows(slide)
    for row_stars, value in zip(rows, star_values):
        full  = int(value)
        half  = (value % 1) >= 0.5
        for j, star in enumerate(row_stars):
            if j < full:
                star.fill.solid()
                star.fill.fore_color.rgb = YELLOW
            elif j == full and half:
                _apply_half_star_fill(star)
            else:
                star.fill.solid()
                star.fill.fore_color.rgb = WHITE


def _is_rating_anchor(shape) -> bool:
    """True if this shape is the rating oval anchor (named 'xx').

    The anchor is the circle in the template.  Its name survives Keynote
    round-trips even when Keynote empties the text and adds a separate TextBox.
    """
    return shape.name.strip().lower() == "xx"


def _is_numeric_rating_text(text: str) -> bool:
    """True if text looks like a filled-in score value."""
    t = text.strip()
    if t.lower() == "xx":
        return True
    return bool(re.fullmatch(r"\d{1,2}([.,]\d{1,2})?", t))


def _is_rating_shape(shape) -> bool:
    """True if this shape holds (or is) the overall rating — used by compatibility check."""
    if _is_rating_anchor(shape):
        return True
    if shape.has_text_frame:
        return _is_numeric_rating_text(shape.text_frame.text)
    return False


def _restore_text_frame(shape) -> None:
    """Re-add a txBody to a shape that Keynote stripped of its text frame."""
    txBody_xml = (
        '<p:txBody xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"'
        ' xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
        '<a:bodyPr anchor="ctr"/>'
        '<a:lstStyle/>'
        '<a:p><a:pPr algn="ctr"/><a:r>'
        '<a:rPr lang="nl-NL" sz="3600" b="1" dirty="0"/>'
        '<a:t></a:t>'
        '</a:r></a:p>'
        '</p:txBody>'
    )
    shape._element.append(etree.fromstring(txBody_xml))


def _find_rating_text_shape(slide):
    """Return the shape whose text should be overwritten with the rating.

    Handles all known layouts produced by Keynote round-trips:

      Case 1 — Original PPTX:
                The 'xx' oval has a text frame containing 'xx' or the score.

      Case 2 — Keynote blank export:
                Keynote wraps the oval in a GroupShape (still named 'xx').
                Inside the group: one circle child + one TextBox child with
                the score text (or empty if never filled).

      Case 3 — Keynote overlay export (filled-then-exported):
                The 'xx' oval is emptied and a separate floating TextBox
                is placed on top of it whose centre aligns with the oval.

    Returns None if nothing useful is found.
    """
    anchor = None
    for shape in slide.shapes:
        if _is_rating_anchor(shape):
            anchor = shape
            break

    if anchor is None:
        return None

    # Case 1 — oval with text frame (original or already filled by this tool)
    if anchor.has_text_frame:
        return anchor

    # Case 2 — Keynote converted oval to a GroupShape named 'xx'
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    if anchor.shape_type == MSO_SHAPE_TYPE.GROUP:
        # Find the TextBox child (the circle child never has useful text)
        for child in anchor.shapes:
            if child.shape_type == MSO_SHAPE_TYPE.TEXT_BOX and child.has_text_frame:
                return child
        # No TextBox child found — return the circle child so we can write into it
        for child in anchor.shapes:
            if child.has_text_frame:
                return child

    # Case 3 — oval emptied, floating TextBox placed on top
    a_cx = anchor.left + anchor.width  // 2
    a_cy = anchor.top  + anchor.height // 2
    margin = anchor.width // 4

    for shape in slide.shapes:
        if shape is anchor:
            continue
        if not shape.has_text_frame:
            continue
        text = shape.text_frame.text.strip()
        if not _is_numeric_rating_text(text):
            continue
        cx = shape.left + shape.width  // 2
        cy = shape.top  + shape.height // 2
        if abs(cx - a_cx) <= margin and abs(cy - a_cy) <= margin:
            return shape

    # Last resort — restore the anchor's text frame if Keynote stripped it
    if not anchor.has_text_frame:
        _restore_text_frame(anchor)
    return anchor


def read_current_star_values(slide) -> list[float]:
    """Return the current rating per row as floats (supports 0.5 for half-stars)."""
    rows = get_star_rows(slide)
    return [sum(_get_star_fill_value(star) for star in row) for row in rows]


def set_rating_text(slide, rating_value: float) -> bool:
    """Write the calculated rating into the rating shape (oval or overlaid TextBox).

    Handles three cases:
      1. Original PPTX  — oval has text frame, text is 'xx' or previous score.
      2. Keynote export — oval is empty, a TextBox sits on top with the score.
      3. Keynote blank  — oval had its text frame deleted entirely by Keynote;
                          we restore the txBody before writing.
    """
    target = _find_rating_text_shape(slide)
    if target is None:
        return False

    # Case 3: Keynote stripped the text frame — restore it first
    if not target.has_text_frame:
        _restore_text_frame(target)

    rating_str = f"{rating_value:.1f}"
    tf = target.text_frame
    first = True
    for para in tf.paragraphs:
        for run in para.runs:
            if first:
                run.text = rating_str
                first = False
            else:
                run.text = ""
    if first:
        para = tf.paragraphs[0]
        para.clear()
        para.add_run().text = rating_str
    return True


def calculate_rating(values: list[int], weights: list[float] | None = None) -> float:
    """Weighted average: sum(v*w) / sum(w), rounded to one decimal."""
    if not values:
        return 0.0
    if not weights or len(weights) != len(values):
        weights = [1.0] * len(values)
    total_weight = sum(weights)
    if total_weight == 0:
        return 0.0
    return round(sum(v * w for v, w in zip(values, weights)) / total_weight, 1)


# ─── Detail slide: comment text ─────────────────────────────────────────────

def get_detail_comment(slide) -> str:
    """Read the scouting comment from TextBox 31 on a detail slide."""
    for shape in slide.shapes:
        if shape.name == "TextBox 31" and shape.has_text_frame:
            return shape.text_frame.text
    return ""


def set_detail_comment(slide, comment: str) -> bool:
    """Write comment text into TextBox 31, preserving its paragraph formatting."""
    for shape in slide.shapes:
        if shape.name == "TextBox 31" and shape.has_text_frame:
            tf = shape.text_frame
            txBody = tf._txBody
            # Keep the first paragraph (contains <a:pPr> with Century Gothic formatting)
            paras = txBody.findall(qn("a:p"))
            for p in paras[1:]:
                txBody.remove(p)
            first_p = paras[0]
            # Strip all runs / line-breaks from it
            for child in list(first_p):
                if child.tag.split("}")[-1] in ("r", "br"):
                    first_p.remove(child)
            lines = comment.split("\n") if comment else [""]
            # Add text to first paragraph
            if lines[0]:
                first_p.append(_make_run(lines[0]))
            # Append extra paragraphs for each additional line (clone formatting)
            for line in lines[1:]:
                new_p = copy.deepcopy(first_p)
                for child in list(new_p):
                    if child.tag.split("}")[-1] in ("r", "br"):
                        new_p.remove(child)
                if line:
                    new_p.append(_make_run(line))
                txBody.append(new_p)
            return True
    return False


def _make_run(text: str):
    """Return an <a:r> lxml element containing the given text."""
    return etree.fromstring(
        f'<a:r xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
        f"<a:t>{text}</a:t></a:r>"
    )


# ─── Detail slide: video ─────────────────────────────────────────────────────

_VIDEO_MIMES = {
    "mp4": "video/mp4",
    "mov": "video/quicktime",
    "avi": "video/avi",
    "wmv": "video/x-ms-wmv",
    "mkv": "video/x-matroska",
    "webm": "video/webm",
}


def _video_mime(filename: str) -> str:
    ext = filename.rsplit(".", 1)[-1].lower() if "." in filename else "mp4"
    return _VIDEO_MIMES.get(ext, "video/mp4")


def extract_first_frame_jpeg(video_bytes: bytes) -> bytes | None:
    """Return JPEG bytes of the first video frame, or None on failure."""
    import os, tempfile
    try:
        import cv2
    except ImportError:
        return None
    with tempfile.NamedTemporaryFile(suffix=".mp4", delete=False) as f:
        f.write(video_bytes)
        tmp = f.name
    try:
        cap = cv2.VideoCapture(tmp)
        ret, frame = cap.read()
        cap.release()
        if not ret or frame is None:
            return None
        ok, buf = cv2.imencode(".jpg", frame, [cv2.IMWRITE_JPEG_QUALITY, 85])
        return bytes(buf) if ok else None
    except Exception:
        return None
    finally:
        try:
            os.unlink(tmp)
        except OSError:
            pass


def _find_detail_placeholder(slide):
    """Return the large placeholder Picture on a detail slide (not the logo)."""
    for shape in slide.shapes:
        if shape.shape_type == 13 and "logo" not in shape.name.lower():  # PICTURE
            return shape
    return None


def get_video_from_slide(slide) -> tuple[bytes, str] | None:
    """Return (bytes, filename) for an embedded video on the slide, or None."""
    for shape in slide.shapes:
        if not isinstance(shape, Movie):
            continue
        # The videoFile element lives under nvPicPr/nvPr with an a: prefix
        nvPr = shape._element.nvPicPr.find(qn("p:nvPr"))
        if nvPr is None:
            continue
        vid = nvPr.find(qn("a:videoFile"))
        if vid is None:
            continue
        r_ns = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
        rId = vid.get(f"{{{r_ns}}}link")
        if rId and rId in slide.part.rels:
            try:
                part = slide.part.rels[rId].target_part
                return part.blob, part.partname.split("/")[-1]
            except Exception:
                pass
    return None


def _remove_movie_shape(slide, shape) -> None:
    """Remove a Movie shape and clean up its orphaned relationships."""
    elem = shape._element
    spTree = slide.shapes._spTree

    # Collect all rIds referenced by this shape element
    r_ns = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
    rIds = set()
    for node in elem.iter():
        for key, val in node.attrib.items():
            if key == f"{{{r_ns}}}embed" or key == f"{{{r_ns}}}link":
                rIds.add(val)

    # Remove the shape element from the slide tree
    spTree.remove(elem)

    # Drop the orphaned relationships so PowerPoint won't ask for repair
    for rId in rIds:
        if rId in slide.part.rels:
            try:
                slide.part.rels.pop(rId)
            except Exception:
                pass


def embed_video_on_slide(
    prs, slide_idx: int, video_bytes: bytes, video_filename: str
) -> bool:
    """Replace the placeholder picture on a detail slide with an embedded video.

    If the slide already has an embedded video it is removed first so we do not
    accumulate duplicate media entries on repeated fills.  The existing Movie's
    geometry is captured *before* removal so re-fills always succeed even after
    the original placeholder Picture has been consumed.
    """
    slide = prs.slides[slide_idx]

    # Capture geometry of any existing Movie BEFORE removing it
    movie_geo = None
    for shape in list(slide.shapes):
        if isinstance(shape, Movie):
            movie_geo = (shape.left, shape.top, shape.width, shape.height)
            _remove_movie_shape(slide, shape)

    # Find placeholder picture and read its geometry
    placeholder = _find_detail_placeholder(slide)
    if placeholder is not None:
        left, top = placeholder.left, placeholder.top
        width, height = placeholder.width, placeholder.height
        slide.shapes._spTree.remove(placeholder._element)
    elif movie_geo is not None:
        # No placeholder (consumed on a previous fill) — reuse Movie's geometry
        left, top, width, height = movie_geo
    else:
        return False

    poster_jpeg = extract_first_frame_jpeg(video_bytes)
    poster_io = io.BytesIO(poster_jpeg) if poster_jpeg else None

    slide.shapes.add_movie(
        io.BytesIO(video_bytes),
        left, top, width, height,
        poster_frame_image=poster_io,
        mime_type=_video_mime(video_filename),
    )
    return True


def _apply_ratings(
    prs,
    template_cfg: dict,
    star_values: list,
    comments: list[str] | None = None,
    video_data: list | None = None,
) -> None:
    """Apply stars, rating text, comments and videos to an open presentation."""
    rating = calculate_rating(star_values, template_cfg.get("weights"))
    main_slide = prs.slides[template_cfg["rating_slide_idx"]]
    color_stars(main_slide, star_values)
    set_rating_text(main_slide, rating)

    detail_idxs = template_cfg.get("detail_slides", [])
    for i, idx in enumerate(detail_idxs):
        slide = prs.slides[idx]
        if i < len(star_values):
            color_stars(slide, [star_values[i]])
        if comments and i < len(comments) and comments[i]:
            set_detail_comment(slide, comments[i])
        if video_data and i < len(video_data) and video_data[i] is not None:
            vb, vname = video_data[i]
            embed_video_on_slide(prs, idx, vb, vname)


# ─── Player info (slide 4 / rating slide) ───────────────────────────────────

# The 9 empty TextBox 11 shapes on the left of the rating slide, sorted by top,
# correspond to these fields in order:
_PLAYER_INFO_FIELDS = [
    "date_of_birth", "city_of_birth", "nationality", "height",
    "preferred_foot", "club", "league", "agency", "agent",
]


def _write_text_shape(shape, text: str) -> None:
    """Write text into a shape's text frame, preserving first-paragraph formatting."""
    if not shape.has_text_frame:
        return
    tf = shape.text_frame
    if tf.paragraphs and tf.paragraphs[0].runs:
        tf.paragraphs[0].runs[0].text = text
        for r in tf.paragraphs[0].runs[1:]:
            r.text = ""
    else:
        para = tf.paragraphs[0]
        para.clear()
        run = para.add_run()
        run.text = text


def _fit_name_to_shape(shape, text: str) -> None:
    """Write name into welcome-slide shape preserving template formatting (64pt)."""
    _write_text_shape(shape, text)


def fill_player_info(prs, template_cfg: dict, player_data: dict) -> None:
    """Fill the player info fields on the rating slide + welcome slide name."""
    rating_slide = prs.slides[template_cfg["rating_slide_idx"]]

    # Fill welcome slide name (TextBox 28 on slide 1), shrinking font if needed
    if len(prs.slides) > 0:
        for shape in prs.slides[0].shapes:
            if shape.name == "TextBox 28" and shape.has_text_frame:
                _fit_name_to_shape(shape, player_data.get("name", ""))
                break

    # Fill the 9 left TextBox 11 fields
    left_fields = []
    for shape in rating_slide.shapes:
        if shape.name == "TextBox 11" and shape.left < 8_000_000:
            left_fields.append((shape.top, shape))
    left_fields.sort()

    for i, (_, shape) in enumerate(left_fields):
        if i < len(_PLAYER_INFO_FIELDS):
            field_name = _PLAYER_INFO_FIELDS[i]
            _write_text_shape(shape, player_data.get(field_name, ""))

    # Fill the Player name textbox
    for shape in rating_slide.shapes:
        if shape.name == "TextBox 37" and shape.has_text_frame:
            _write_text_shape(shape, player_data.get("name", ""))
            break


# ─── Player stats (right side of rating slide) ────────────────────────────

# The right-side TextBox 11 shapes (left > 8M EMU) are stat value fields.
# They are arranged in 4 rows x 2 columns:
#   Row 1 (top ~1.1M):  season_matches  |  career_matches
#   Row 2 (top ~1.8M):  season_minutes  |  career_minutes
#   Row 3 (top ~2.6M):  season_goals    |  career_goals
#   Row 4 (top ~3.4M):  season_assists  |  career_assists
# Within each row, the left column is ~11.2M and right is ~13.6M.
_STAT_FIELD_ORDER = [
    # (row, col) -> field_name — sorted by top then by left
    "season_matches", "career_matches",
    "season_minutes", "career_minutes",
    "season_goals", "career_goals",
    "season_assists", "career_assists",
]


def fill_player_stats(prs, template_cfg: dict, tm_stats: dict) -> None:
    """Fill the season/career stat fields on the rating slide."""
    rating_slide = prs.slides[template_cfg["rating_slide_idx"]]

    # Collect right-side TextBox 11 shapes
    right_fields = []
    for shape in rating_slide.shapes:
        if shape.name == "TextBox 11" and shape.left >= 8_000_000:
            right_fields.append((shape.top, shape.left, shape))

    if not right_fields:
        return

    # Sort by top first, then left within each row
    right_fields.sort(key=lambda x: (x[0], x[1]))

    # Group into rows (shapes within 200K EMU of each other vertically)
    rows = []
    current_row = [right_fields[0]]
    for item in right_fields[1:]:
        if abs(item[0] - current_row[0][0]) < 200_000:
            current_row.append(item)
        else:
            rows.append(sorted(current_row, key=lambda x: x[1]))
            current_row = [item]
    rows.append(sorted(current_row, key=lambda x: x[1]))

    # Flatten rows into field order and fill
    idx = 0
    for row in rows:
        for _, _, shape in row:
            if idx < len(_STAT_FIELD_ORDER):
                field = _STAT_FIELD_ORDER[idx]
                value = tm_stats.get(field, 0)
                _write_text_shape(shape, str(value) if value else "0")
                idx += 1


# ─── Transfer Details (rating slide, bottom-right "xxxx" column) ──────────

# The 5 textboxes named "xxxx" on the rating slide, sorted top→bottom:
#   0: End of contract     e.g. "Jun 2026"
#   1: Transfer value      e.g. "€750K"
#   2: Prediction year 1   e.g. "Top KKD"
#   3: Prediction year 2   e.g. "Eredivisie"
#   4: Next step           e.g. "FC Den Bosch"
TRANSFER_FIELD_ORDER = [
    "end_of_contract", "transfer_value", "prediction_year_1",
    "prediction_year_2", "next_step",
]


def fill_transfer_details(prs, template_cfg: dict, transfer_details: dict) -> None:
    """Fill the 5 'xxxx' placeholders on the rating slide with transfer details."""
    if not transfer_details:
        return
    rating_slide = prs.slides[template_cfg["rating_slide_idx"]]

    xxxx_shapes = [
        s for s in rating_slide.shapes
        if s.name == "xxxx" and s.has_text_frame
    ]
    xxxx_shapes.sort(key=lambda s: s.top or 0)

    for i, shape in enumerate(xxxx_shapes[: len(TRANSFER_FIELD_ORDER)]):
        field = TRANSFER_FIELD_ORDER[i]
        val = transfer_details.get(field)
        if val:
            _write_text_shape(shape, str(val))


# ─── Scouting dates (rating slide, top-right "TextBox 23") ────────────────

def fill_scouting_dates(prs, template_cfg: dict, scouting_dates: list) -> None:
    """Fill 'TextBox 23' on the rating slide with one line per scouting entry.

    Each entry is a dict {"date": "DD/MM/YYYY", "type": "Game"|"Training"}.
    """
    if not scouting_dates:
        return
    rating_slide = prs.slides[template_cfg["rating_slide_idx"]]

    target = None
    for shape in rating_slide.shapes:
        if shape.name == "TextBox 23" and shape.has_text_frame:
            target = shape
            break
    if target is None:
        return

    lines = []
    for entry in scouting_dates:
        d = (entry or {}).get("date", "").strip()
        ttype = (entry or {}).get("type", "").strip()
        if d and ttype:
            lines.append(f"{d}: {ttype}")
        elif d:
            lines.append(d)
    if lines:
        _write_text_shape(target, "\n".join(lines))


# ─── Physical data ─────────────────────────────────────────────────────────

def _write_multiline_text_shape(shape, lines: list) -> None:
    """Write one line per paragraph into a shape, preserving each paragraph's
    existing first-run formatting (font size / colour / family)."""
    if not shape.has_text_frame:
        return
    tf = shape.text_frame
    paragraphs = tf.paragraphs
    for i, line in enumerate(lines):
        text = "" if line is None else str(line)
        if i < len(paragraphs):
            p = paragraphs[i]
            if p.runs:
                p.runs[0].text = text
                for r in p.runs[1:]:
                    r.text = ""
            else:
                p.clear()
                r = p.add_run()
                r.text = text


def fill_physical_data(prs, template_cfg: dict, physical_data: dict) -> None:
    """Fill the physical-data placeholders (bottom-left stack) on the rating slide.

    ``TextBox 31`` (top~12.64″, left~2.85″) — 4 lines:
    Total Distance / HI-runs / Sprints / Top speed.
    """
    if not physical_data:
        return
    rating_slide = prs.slides[template_cfg["rating_slide_idx"]]

    def _fmt_distance(v):
        try:
            return f"{float(v) / 1000.0:.2f} km"
        except (TypeError, ValueError):
            return ""

    def _fmt_int(v):
        try:
            return str(int(round(float(v))))
        except (TypeError, ValueError):
            return ""

    def _fmt_speed(v):
        if v is None:
            return ""
        s = str(v).strip()
        if not s:
            return ""
        low = s.lower().replace(" ", "")
        if low.endswith("km/h") or low.endswith("kmh") or low.endswith("kph"):
            return s
        try:
            return f"{float(s):.1f} km/h"
        except ValueError:
            return s

    total_distance = physical_data.get("total_distance")
    hi_runs = physical_data.get("hi_runs")
    sprints = physical_data.get("sprint_efforts")
    top_speed = physical_data.get("top_speed")

    EMU = 914400
    for shape in rating_slide.shapes:
        if (shape.name == "TextBox 31"
                and shape.has_text_frame
                and (shape.top or 0) > int(11.5 * EMU)
                and (shape.left or 0) < int(5 * EMU)):
            _write_multiline_text_shape(shape, [
                _fmt_distance(total_distance),
                _fmt_int(hi_runs),
                _fmt_int(sprints),
                _fmt_speed(top_speed),
            ])
            break


def fill_availability(prs, template_cfg: dict, availability_pct) -> None:
    """Fill the big ``xx%`` box (PLAYER AVAILABILITY, top~5.60″, left~13.88″)
    with the scraped Transfermarkt availability percentage."""
    if availability_pct is None:
        return
    try:
        pct_text = f"{float(availability_pct):.0f}%"
    except (TypeError, ValueError):
        return
    rating_slide = prs.slides[template_cfg["rating_slide_idx"]]
    for shape in rating_slide.shapes:
        if shape.name == "xx%" and shape.has_text_frame:
            _write_text_shape(shape, pct_text)
            break


def fill_player_photo(
    prs,
    template_cfg: dict,
    full_photo: bytes | None = None,
    circular_photo: bytes | None = None,
) -> None:
    """Place player photos in the presentation.

    * **full_photo** → welcome slide (slide 0), left side
    * **circular_photo** → rating slide, replacing the 'Browse' placeholder
    """
    # ── Full image on welcome slide (slide 0) ──────────────────────────────
    if full_photo:
        from PIL import Image as _PILImage
        slide0 = prs.slides[0]

        # Find reference shapes to position the image:
        #   • left bound  = left edge of name bar (Rechthoek)
        #   • right bound = left edge of WELKOM/WELCOME text
        #   • bottom      = top of the Rechthoek (name-bar background)

        welkom_left = int(10.09 * 914400)    # fallback FC Den Bosch
        bar_top = int(8.18 * 914400)         # fallback
        bar_left = int(4.97 * 914400)        # fallback

        # Search layout first, then master as fallback
        _welkom_found = False
        _bar_found = False
        _search_pools = [slide0.slide_layout.shapes]
        try:
            _search_pools.append(slide0.slide_layout.slide_master.shapes)
        except Exception:
            pass
        for pool in _search_pools:
            for shape in pool:
                nm = shape.name or ""
                txt = shape.text if hasattr(shape, "text") else ""
                if not _welkom_found and ("WELKOM" in txt or "WELCOME" in txt):
                    welkom_left = shape.left
                    _welkom_found = True
                if not _bar_found and nm == "Rechthoek":
                    bar_top = shape.top
                    bar_left = shape.left
                    _bar_found = True

        # Target box: from Rechthoek left to WELKOM left, top y=0 to bar top
        box_left = bar_left
        box_right = welkom_left
        box_w = box_right - box_left
        box_top = 0
        box_h = bar_top - box_top

        _img = _PILImage.open(io.BytesIO(full_photo))
        iw, ih = _img.size
        scale = min(box_w / iw, box_h / ih)
        pic_w = int(iw * scale)
        pic_h = int(ih * scale)
        # Centre horizontally within our new left-side box, anchor bottom to bar top
        pic_left = box_left + (box_w - pic_w) // 2
        pic_top = bar_top - pic_h

        # Remove any previously added player photo (re-fill safe)
        _spTree0 = slide0.shapes._spTree
        for shape in list(slide0.shapes):
            if shape.name == "player_photo_welcome":
                _spTree0.remove(shape._element)

        img_stream = io.BytesIO(full_photo)
        pic = slide0.shapes.add_picture(img_stream, pic_left, pic_top, pic_w, pic_h)
        pic.name = "player_photo_welcome"

    # ── Circular crop on rating slide ──────────────────────────────────────
    photo_for_rating = circular_photo or full_photo
    if photo_for_rating:
        rating_slide = prs.slides[template_cfg["rating_slide_idx"]]

        # Find placement: "Browse" placeholder (first fill) or previously placed photo (re-fill)
        target_shape = None
        for shape in rating_slide.shapes:
            if shape.name.startswith("Browse"):
                target_shape = shape
                break
        if target_shape is None:
            for shape in rating_slide.shapes:
                if shape.name == "player_photo_rating":
                    target_shape = shape
                    break

        if target_shape is not None:
            left = target_shape.left
            top = target_shape.top
            width = target_shape.width
            height = target_shape.height
            sp = target_shape._element
            sp.getparent().remove(sp)
            img_stream = io.BytesIO(photo_for_rating)
            pic = rating_slide.shapes.add_picture(img_stream, left, top, width, height)
            pic.name = "player_photo_rating"

def _today_ddmmyyyy() -> str:
    import datetime as _dt
    return _dt.date.today().strftime("%d-%m-%Y")


def extract_report_date(pptx_bytes: bytes, rating_slide_idx: int = 3) -> str | None:
    """Return the current 'DATE:' value from an existing pptx, or None if the
    field is still the placeholder (DD-MM-JJJJ) or cannot be found."""
    try:
        prs = Presentation(io.BytesIO(pptx_bytes))
    except Exception:
        return None
    if rating_slide_idx >= len(prs.slides):
        return None
    for shape in prs.slides[rating_slide_idx].shapes:
        if not shape.has_text_frame:
            continue
        txt = (shape.text_frame.text or "").strip()
        if txt.upper().startswith("DATE:"):
            val = txt.split(":", 1)[1].strip()
            if not val or "JJJJ" in val.upper() or "YYYY" in val.upper():
                return None
            return val
    return None


def fill_report_date(prs, template_cfg: dict, date_str: str) -> None:
    """Replace the 'DATE: DD-MM-JJJJ' placeholder on the rating slide with the
    given date string (formatted as DD-MM-YYYY by the caller)."""
    if not date_str:
        return
    rating_slide = prs.slides[template_cfg["rating_slide_idx"]]
    for shape in rating_slide.shapes:
        if not shape.has_text_frame:
            continue
        txt = shape.text_frame.text or ""
        if "DD-MM-JJJJ" in txt or "DATE:" in txt.upper():
            _write_text_shape(shape, f"DATE: {date_str}")
            return


def fill_template(
    template_cfg: dict,
    star_values: list,
    comments: list[str] | None = None,
    video_data: list | None = None,
    player_data: dict | None = None,
    tm_stats: dict | None = None,
    player_photo: bytes | None = None,
    player_photo_circular: bytes | None = None,
    physical_data: dict | None = None,
    transfer_details: dict | None = None,
    scouting_dates: list | None = None,
    report_date: str | None = None,
) -> io.BytesIO:
    """Fill a blank template file and return the result as BytesIO."""
    prs = Presentation(template_cfg["file"])
    if player_data:
        fill_player_info(prs, template_cfg, player_data)
    if tm_stats:
        fill_player_stats(prs, template_cfg, tm_stats)
        fill_availability(prs, template_cfg, tm_stats.get("availability_pct"))
    if physical_data:
        fill_physical_data(prs, template_cfg, physical_data)
    if transfer_details:
        fill_transfer_details(prs, template_cfg, transfer_details)
    if scouting_dates:
        fill_scouting_dates(prs, template_cfg, scouting_dates)
    fill_report_date(prs, template_cfg, report_date or _today_ddmmyyyy())
    if player_photo or player_photo_circular:
        fill_player_photo(prs, template_cfg, full_photo=player_photo, circular_photo=player_photo_circular)
    _apply_ratings(prs, template_cfg, star_values, comments, video_data)
    output = io.BytesIO()
    prs.save(output)
    output.seek(0)
    return output


def fill_from_bytes(
    file_bytes: bytes,
    template_cfg: dict,
    star_values: list,
    comments: list[str] | None = None,
    video_data: list | None = None,
    player_data: dict | None = None,
    tm_stats: dict | None = None,
    player_photo: bytes | None = None,
    player_photo_circular: bytes | None = None,
    physical_data: dict | None = None,
    transfer_details: dict | None = None,
    scouting_dates: list | None = None,
    report_date: str | None = None,
) -> io.BytesIO:
    """Fill an uploaded PPTX (raw bytes) and return the result as BytesIO."""
    prs = Presentation(io.BytesIO(file_bytes))
    if player_data:
        fill_player_info(prs, template_cfg, player_data)
    if tm_stats:
        fill_player_stats(prs, template_cfg, tm_stats)
        fill_availability(prs, template_cfg, tm_stats.get("availability_pct"))
    if physical_data:
        fill_physical_data(prs, template_cfg, physical_data)
    if transfer_details:
        fill_transfer_details(prs, template_cfg, transfer_details)
    if scouting_dates:
        fill_scouting_dates(prs, template_cfg, scouting_dates)
    fill_report_date(prs, template_cfg, report_date or _today_ddmmyyyy())
    if player_photo or player_photo_circular:
        fill_player_photo(prs, template_cfg, full_photo=player_photo, circular_photo=player_photo_circular)
    _apply_ratings(prs, template_cfg, star_values, comments, video_data)
    output = io.BytesIO()
    prs.save(output)
    output.seek(0)
    return output


# ─── Slide image export (for in-app preview) ─────────────────────────────

# Module-level diagnostic + cache for the preview render.
_last_preview_error: str = ""
_preview_cache: "OrderedDict[str, tuple[str, bytes]]"  # populated lazily
_PREVIEW_CACHE_MAX = 16
_LO_PROFILE_DIR: str | None = None  # reused LibreOffice user profile (faster warm starts)

try:  # keep OrderedDict import local-friendly but resolved at module load
    from collections import OrderedDict as _OD
    _preview_cache = _OD()
except Exception:
    _preview_cache = {}  # type: ignore


def get_last_preview_error() -> str:
    """Error message from the most recent render_slide_preview call (or '')."""
    return _last_preview_error


def _cache_key(pptx_bytes: bytes, slide_index: int, width: int, kind_hint: str = "") -> str:
    import hashlib
    h = hashlib.sha256(pptx_bytes).hexdigest()[:16]
    return f"{h}-{slide_index}-{width}-{kind_hint}"


def _cache_get(key: str) -> tuple[str, bytes] | None:
    hit = _preview_cache.get(key)
    if hit is not None:
        try:
            _preview_cache.move_to_end(key)  # LRU bump
        except Exception:
            pass
    return hit


def _cache_put(key: str, value: tuple[str, bytes]) -> None:
    _preview_cache[key] = value
    while len(_preview_cache) > _PREVIEW_CACHE_MAX:
        try:
            _preview_cache.popitem(last=False)  # drop oldest
        except Exception:
            _preview_cache.pop(next(iter(_preview_cache)))


def _lo_profile() -> str:
    """Persistent LibreOffice user profile dir — keeps warm-start caches."""
    global _LO_PROFILE_DIR
    if _LO_PROFILE_DIR and os.path.isdir(_LO_PROFILE_DIR):
        return _LO_PROFILE_DIR
    import os as _os
    import tempfile as _tf
    _LO_PROFILE_DIR = _tf.mkdtemp(prefix="lo_profile_")
    return _LO_PROFILE_DIR


def _pdf_page_to_png_via_pdftoppm(pdf_path: str, page_1based: int, out_dir: str,
                                  width: int) -> bytes | None:
    """Use poppler's pdftoppm (if installed) to rasterise a single PDF page to
    PNG. Returns PNG bytes, or None if unavailable/failed."""
    import shutil
    import subprocess
    import os as _os
    import glob

    tool = shutil.which("pdftoppm")
    if not tool:
        return None
    out_base = _os.path.join(out_dir, "slide")
    try:
        # Use -scale-to-x for exact width; height auto.
        proc = subprocess.run(
            [tool, "-png", "-f", str(page_1based), "-l", str(page_1based),
             "-scale-to-x", str(width), "-scale-to-y", "-1",
             pdf_path, out_base],
            capture_output=True, timeout=30,
        )
        if proc.returncode != 0:
            return None
    except Exception:
        return None
    # pdftoppm names output like slide-1.png (single-digit) or slide-01.png.
    matches = sorted(glob.glob(out_base + "*.png"))
    if not matches:
        return None
    try:
        with open(matches[0], "rb") as f:
            return f.read()
    except Exception:
        return None


def _render_via_libreoffice(
    src_path: str,
    tmpdir: str,
    slide_index: int,
    width: int = 1280,
) -> tuple[tuple[str, bytes] | None, str | None]:
    """Convert pptx → pdf via headless LibreOffice. Returns ((kind, bytes), None)
    on success or (None, error_msg) on failure.

    Post-processing priority:
      1. pdftoppm → single PNG of target page (fast display in browser)
      2. pypdf    → single-page PDF of target page
      3. fall back to whole-deck PDF
    """
    import os
    import shutil
    import subprocess

    _ensure_font_aliases()

    soffice = shutil.which("soffice") or shutil.which("libreoffice")
    if not soffice:
        return None, (
            "LibreOffice not found on PATH. Install with: "
            "apt-get install -y libreoffice-impress  (Linux) "
            "or brew install --cask libreoffice (macOS)."
        )

    profile = _lo_profile()
    try:
        proc = subprocess.run(
            [
                soffice,
                f"-env:UserInstallation=file://{profile}",
                "--headless", "--norestore", "--nologo",
                "--nolockcheck", "--nofirststartwizard",
                "--convert-to", "pdf", "--outdir", tmpdir, src_path,
            ],
            capture_output=True,
            timeout=180,
        )
    except Exception as exc:
        return None, f"LibreOffice subprocess failed: {exc}"

    if proc.returncode != 0:
        err = (proc.stderr or b"").decode("utf-8", "ignore")[:400]
        return None, f"LibreOffice exited {proc.returncode}: {err}"

    base = os.path.splitext(os.path.basename(src_path))[0]
    pdf_path = os.path.join(tmpdir, base + ".pdf")
    if not os.path.exists(pdf_path) or os.path.getsize(pdf_path) == 0:
        return None, "LibreOffice produced no PDF"

    # Strategy 1: fast PNG of just the target slide via pdftoppm
    png_bytes = _pdf_page_to_png_via_pdftoppm(pdf_path, slide_index + 1, tmpdir, width)
    if png_bytes:
        return ("png", png_bytes), None

    # Strategy 2: single-page PDF via pypdf
    try:
        from pypdf import PdfReader, PdfWriter  # type: ignore
        reader = PdfReader(pdf_path)
        if 0 <= slide_index < len(reader.pages):
            writer = PdfWriter()
            writer.add_page(reader.pages[slide_index])
            import io
            buf = io.BytesIO()
            writer.write(buf)
            return ("pdf", buf.getvalue()), None
    except Exception:
        pass

    # Strategy 3: whole-deck PDF
    with open(pdf_path, "rb") as f:
        return ("pdf", f.read()), None


def render_slide_preview(
    pptx_bytes: bytes,
    slide_index: int,
    width: int = 1280,
) -> tuple[str, bytes] | None:
    """Render one slide to a preview. Returns ("png", bytes) or ("pdf", bytes).

    Results are cached by SHA-256 of pptx_bytes, so re-previewing the same
    content is instant. Strategy order:
      1. PowerPoint COM (Windows, Office installed) → PNG then PDF
      2. LibreOffice headless (Linux/Mac/Windows if installed) → PNG via
         pdftoppm, else single-page PDF via pypdf, else whole-deck PDF
      3. Returns None; call get_last_preview_error() for the reason
    """
    global _last_preview_error
    _last_preview_error = ""
    errors: list[str] = []

    cache_key = _cache_key(pptx_bytes, slide_index, width)
    cached = _cache_get(cache_key)
    if cached is not None:
        return cached

    import tempfile

    tmpdir = tempfile.mkdtemp(prefix="ppt_preview_")
    src_path = os.path.join(tmpdir, "src.pptx")
    png_path = os.path.join(tmpdir, f"slide_{slide_index + 1}.png")
    pdf_path = os.path.join(tmpdir, "preview.pdf")

    try:
        with open(src_path, "wb") as f:
            f.write(pptx_bytes)
    except Exception as exc:
        _last_preview_error = f"Could not write temp pptx: {exc}"
        _cleanup_tmpdir(tmpdir)
        return None

    # ── Strategy A: PowerPoint COM (Windows-only) ─────────────────────
    try:
        import pythoncom  # type: ignore
        import win32com.client  # type: ignore  # noqa: F401
        _pywin32_ok = True
    except Exception as exc:
        errors.append(f"pywin32 not available: {exc}")
        _pywin32_ok = False

    if _pywin32_ok:
        result = _render_via_powerpoint_com(
            src_path, png_path, pdf_path, slide_index, width, errors,
        )
        if result is not None:
            _cache_put(cache_key, result)
            _cleanup_tmpdir(tmpdir)
            return result

    # ── Strategy B: LibreOffice headless (any OS if installed) ────────
    lo_result, lo_err = _render_via_libreoffice(src_path, tmpdir, slide_index, width)
    if lo_result is not None:
        _cache_put(cache_key, lo_result)
        _cleanup_tmpdir(tmpdir)
        return lo_result
    if lo_err:
        errors.append(lo_err)

    _last_preview_error = "; ".join(errors) or "unknown failure"
    _cleanup_tmpdir(tmpdir)
    return None


_FONT_ALIASES_WRITTEN = False


def _ensure_font_aliases() -> None:
    """Install a user fontconfig file that aliases the template's Mac/Adobe fonts
    (Helvetica Neue, Avenir Next) to closest-metric Linux substitutes so the
    LibreOffice preview doesn't fall back to DejaVu (wider glyphs → oversized
    white letters in the rendered image). Safe to call many times.
    """
    global _FONT_ALIASES_WRITTEN
    if _FONT_ALIASES_WRITTEN:
        return
    try:
        home = os.path.expanduser("~")
        cfg_dir = os.path.join(home, ".config", "fontconfig")
        os.makedirs(cfg_dir, exist_ok=True)
        cfg_path = os.path.join(cfg_dir, "fonts.conf")
        fonts_conf = """<?xml version="1.0"?>
<!DOCTYPE fontconfig SYSTEM "fonts.dtd">
<fontconfig>
  <alias binding="strong"><family>Helvetica Neue</family><prefer><family>Nimbus Sans</family><family>Liberation Sans</family><family>Arial</family></prefer></alias>
  <alias binding="strong"><family>Helvetica Neue Medium</family><prefer><family>Nimbus Sans</family><family>Liberation Sans</family><family>Arial</family></prefer></alias>
  <alias binding="strong"><family>Helvetica</family><prefer><family>Nimbus Sans</family><family>Liberation Sans</family><family>Arial</family></prefer></alias>
  <alias binding="strong"><family>Avenir Next Condensed</family><prefer><family>Nimbus Sans Narrow</family><family>Liberation Sans Narrow</family><family>Nimbus Sans</family></prefer></alias>
  <alias binding="strong"><family>Avenir Next Condensed Regular</family><prefer><family>Nimbus Sans Narrow</family><family>Liberation Sans Narrow</family><family>Nimbus Sans</family></prefer></alias>
  <alias binding="strong"><family>Avenir Next</family><prefer><family>Nimbus Sans</family><family>Liberation Sans</family></prefer></alias>
  <alias binding="strong"><family>Avenir</family><prefer><family>Nimbus Sans</family><family>Liberation Sans</family></prefer></alias>
</fontconfig>
"""
        with open(cfg_path, "w", encoding="utf-8") as f:
            f.write(fonts_conf)
        _FONT_ALIASES_WRITTEN = True
    except Exception:
        pass


def warm_up_preview_engine() -> None:
    """Pre-start LibreOffice in the background so the first real preview is fast.

    Non-blocking: spawns a detached soffice process that primes the user profile
    and font cache, then exits. Safe to call many times.
    """
    try:
        import shutil
        import subprocess
        _ensure_font_aliases()
        soffice = shutil.which("soffice") or shutil.which("libreoffice")
        if not soffice:
            return
        profile = _lo_profile()
        # --terminate_after_init warms caches then quits immediately
        subprocess.Popen(
            [
                soffice,
                f"-env:UserInstallation=file://{profile}",
                "--headless", "--norestore", "--nologo",
                "--nolockcheck", "--nofirststartwizard",
                "--terminate_after_init",
            ],
            stdout=subprocess.DEVNULL,
            stderr=subprocess.DEVNULL,
        )
    except Exception:
        pass


def _cleanup_tmpdir(tmpdir: str) -> None:
    try:
        import shutil
        shutil.rmtree(tmpdir, ignore_errors=True)
    except Exception:
        pass


def _render_via_powerpoint_com(
    src_path: str,
    png_path: str,
    pdf_path: str,
    slide_index: int,
    width: int,
    errors: list[str],
) -> tuple[str, bytes] | None:
    """Render via PowerPoint COM (Windows only). Appends diagnostics to `errors`."""
    import os
    import pythoncom  # type: ignore
    import win32com.client  # type: ignore

    pythoncom.CoInitialize()
    ppt_app = None
    pres = None
    try:
        try:
            ppt_app = win32com.client.Dispatch("PowerPoint.Application")
        except Exception as exc:
            errors.append(f"PowerPoint not available (COM Dispatch failed: {exc})")
            return None

        # PowerPoint requires a visible window for many COM operations.
        # Keeping Visible=1 + WindowState=2 (minimised) is the usual workaround.
        try:
            ppt_app.Visible = 1
            try:
                ppt_app.WindowState = 2  # ppWindowMinimized
            except Exception:
                pass
        except Exception as exc:
            errors.append(f"PowerPoint.Visible=1 failed: {exc}")

        try:
            pres = ppt_app.Presentations.Open(
                src_path, ReadOnly=True, Untitled=False, WithWindow=False
            )
        except Exception as exc:
            errors.append(f"Presentations.Open failed: {exc}")
            return None

        # Strategy 1 — single-slide PNG export
        try:
            slide_count = int(pres.Slides.Count)
            if slide_index + 1 <= slide_count:
                slide = pres.Slides(slide_index + 1)
                try:
                    height = int(width * (pres.PageSetup.SlideHeight / pres.PageSetup.SlideWidth))
                except Exception:
                    height = int(width * 9 / 16)
                slide.Export(png_path, "PNG", width, height)
                if os.path.exists(png_path) and os.path.getsize(png_path) > 0:
                    with open(png_path, "rb") as f:
                        return ("png", f.read())
                errors.append("Slide.Export produced no file")
            else:
                errors.append(
                    f"slide_index {slide_index} out of range (deck has {slide_count} slides)"
                )
        except Exception as exc:
            errors.append(f"Slide.Export PNG failed: {exc}")

        # Strategy 2 — whole-deck PDF export (ppSaveAsPDF = 32)
        try:
            pres.SaveAs(pdf_path, 32)
            if os.path.exists(pdf_path) and os.path.getsize(pdf_path) > 0:
                with open(pdf_path, "rb") as f:
                    return ("pdf", f.read())
            errors.append("Presentation.SaveAs PDF produced no file")
        except Exception as exc:
            errors.append(f"Presentation.SaveAs PDF failed: {exc}")

        return None
    except Exception as exc:
        errors.append(f"Unexpected: {exc}")
        return None
    finally:
        try:
            if pres is not None:
                pres.Close()
        except Exception:
            pass
        try:
            if ppt_app is not None:
                ppt_app.Quit()
        except Exception:
            pass
        try:
            pythoncom.CoUninitialize()
        except Exception:
            pass


def render_slide_as_image(pptx_bytes: bytes, slide_index: int, width: int = 1280) -> bytes | None:
    """Back-compat wrapper — returns PNG bytes only (or None)."""
    result = render_slide_preview(pptx_bytes, slide_index, width)
    if result and result[0] == "png":
        return result[1]
    return None


# ─── Template detection & compatibility ─────────────────────────────────────

def detect_template_name(slide) -> tuple[str | None, str | None]:
    """Return (position_name, language) for the best-matching template, or (None, None)."""
    slide_text = " ".join(
        shape.text_frame.text.lower()
        for shape in slide.shapes
        if shape.has_text_frame
    )
    best_name: str | None = None
    best_lang: str | None = None
    best_count = 0
    for name, cfg in TEMPLATES.items():
        for lang_key, lang_label in [("variables_nl", "NL"), ("variables_eng", "ENG")]:
            matches = sum(1 for v in cfg[lang_key] if v.lower() in slide_text)
            if matches > best_count:
                best_count = matches
                best_name = name
                best_lang = lang_label
    if best_count >= 3:
        return best_name, best_lang
    return None, None


def check_template_compatibility(file_obj) -> dict:
    """Inspect an uploaded PPTX and return a compatibility report.

    Keys: compatible, star_count, row_count, slide_idx,
          has_rating_placeholder, matched_template_name,
          current_star_values, issues
    """
    result = {
        "compatible": False,
        "star_count": 0,
        "row_count": 0,
        "slide_idx": None,
        "has_rating_placeholder": False,
        "matched_template_name": None,
        "current_star_values": [],
        "issues": [],
    }

    try:
        prs = Presentation(file_obj)
    except Exception as exc:
        result["issues"].append(f"Could not open file: {exc}")
        return result

    # Find the slide with the most star rows
    best_slide_idx = None
    best_rows: list = []
    for idx, slide in enumerate(prs.slides):
        rows = get_star_rows(slide)
        if len(rows) > len(best_rows):
            best_rows = rows
            best_slide_idx = idx

    if not best_rows:
        result["issues"].append("No star shapes found in any slide.")
        return result

    slide = prs.slides[best_slide_idx]
    result["slide_idx"]  = best_slide_idx
    result["row_count"]  = len(best_rows)
    result["star_count"] = sum(len(r) for r in best_rows)

    bad_rows = [i + 1 for i, r in enumerate(best_rows) if len(r) != 10]
    if bad_rows:
        result["issues"].append(
            f"Expected 10 stars per row; rows {bad_rows} have a different count."
        )

    has_placeholder = _find_rating_text_shape(slide) is not None
    result["has_rating_placeholder"] = has_placeholder
    if not has_placeholder:
        result["issues"].append(
            "Rating circle not found. "
            "The file may be too heavily restructured to fill automatically."
        )

    result["current_star_values"] = read_current_star_values(slide)

    matched_name, matched_lang = detect_template_name(slide)
    result["matched_template_name"] = matched_name
    result["matched_language"]      = matched_lang
    # Detect club from slide position: Pro Vercelli uses rating_slide_idx 1
    result["matched_club"] = "Pro Vercelli" if best_slide_idx == 1 else "FC Den Bosch"
    result["compatible"]   = len(result["issues"]) == 0

    # Extract per-detail-slide data (comments + videos) when a template is matched
    if matched_name and matched_name in TEMPLATES:
        club = result["matched_club"]
        lang = matched_lang or "ENG"
        variant_key = (club, lang)
        if variant_key in TEMPLATES[matched_name]["variants"]:
            detail_idxs = TEMPLATES[matched_name]["variants"][variant_key]["detail_slides"]
        else:
            # Fallback: try any variant for this club
            detail_idxs = []
            for k, v in TEMPLATES[matched_name]["variants"].items():
                if k[0] == club:
                    detail_idxs = v["detail_slides"]
                    break
        comments, videos = [], []
        for idx in detail_idxs:
            if idx < len(prs.slides):
                ds = prs.slides[idx]
                comments.append(get_detail_comment(ds))
                videos.append(get_video_from_slide(ds))
            else:
                comments.append("")
                videos.append(None)
        result["current_comments"] = comments
        result["current_videos"]   = videos
    else:
        result["current_comments"] = []
        result["current_videos"]   = []

    return result